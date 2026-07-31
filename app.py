"""
学习助手 Pro - DDL 管理与学习资料库系统（v2 新增个性化复习方案）
========================================
新增功能：
  1. 首次进入引导用户配置每日学习时间段，生成个人学习画像
  2. 基于学习画像 + 当前任务，一键生成可复制的私人复习方案 Prompt
"""
import streamlit as st
import pandas as pd
import requests
import json
import os
from datetime import datetime, timedelta
import base64
import plotly.express as px

# ============================================================
# 一、依赖库加载（带容错）
# ============================================================
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from icalendar import Calendar, Event
except ImportError:
    Calendar = None
    Event = None

try:
    from dateutil import parser as date_parser
except ImportError:
    date_parser = None

# ============================================================
# 二、页面配置
# ============================================================
st.set_page_config(
    page_title="学习助手 Pro",
    layout="centered",
    initial_sidebar_state="collapsed"
)

st.markdown("""
<style>
    .stApp { max-width: 100%; padding: 0.5rem; }
    .stDataFrame { font-size: 12px; }
    .stButton button { width: 100%; margin: 0.2rem 0; }
    @media (max-width: 600px) {
        .row-widget.stColumns { flex-direction: column !important; }
    }
    .cal-grid { display: grid; grid-template-columns: repeat(7, 1fr); gap: 2px; }
    .cal-cell { min-height: 60px; background: #f9f9f9; border-radius: 4px;
                padding: 2px; border-top: 3px solid #ddd; overflow: hidden;
                font-size: 11px; }
    .cal-cell .date { font-weight: bold; font-size: 13px; }
    .cal-weekday { text-align: center; font-weight: bold; color: #888;
                   font-size: 13px; padding: 4px 0; }
    /* ===== 【新增】学习画像卡片样式 ===== */
    .profile-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        border-radius: 12px; padding: 1.2rem; color: white;
        margin-bottom: 1rem;
    }
    .profile-card h4 { margin: 0 0 0.5rem 0; color: #fff; }
    .profile-card .tag {
        display: inline-block; background: rgba(255,255,255,0.2);
        border-radius: 20px; padding: 4px 12px; margin: 3px 4px;
        font-size: 13px;
    }
    .profile-card .info { font-size: 14px; opacity: 0.9; margin-top: 8px; }
</style>
""", unsafe_allow_html=True)

# ============================================================
# 三、数据文件路径
# ============================================================
DDL_FILE = "deadlines.csv"
LIBRARY_FILE = "library.csv"
CATEGORIES_FILE = "categories.csv"
PROFILE_FILE = "study_profile.csv"  # ===== 【新增】学习画像文件 =====

# ============================================================
# 四、数据初始化
# ============================================================

# --- DDL 数据 ---
if "df" not in st.session_state:
    if os.path.exists(DDL_FILE):
        st.session_state.df = pd.read_csv(DDL_FILE)
        for col in ["重复", "状态"]:
            if col not in st.session_state.df.columns:
                st.session_state.df[col] = "无" if col == "重复" else "未完成"
    else:
        st.session_state.df = pd.DataFrame(columns=[
            "课程/科目", "截止日期", "描述", "标签", "重复", "状态", "添加时间"
        ])

# --- 资料库 ---
if "library" not in st.session_state:
    if os.path.exists(LIBRARY_FILE):
        st.session_state.library = pd.read_csv(LIBRARY_FILE)
    else:
        st.session_state.library = pd.DataFrame(columns=[
            "文件名", "分类", "摘要", "上传时间", "内容"
        ])

# --- 分类 ---
if "categories" not in st.session_state:
    if os.path.exists(CATEGORIES_FILE):
        st.session_state.categories = pd.read_csv(CATEGORIES_FILE)["分类"].tolist()
    else:
        st.session_state.categories = ["未分类"]

# ===== 【新增】学习画像初始化 =====
if "profile" not in st.session_state:
    if os.path.exists(PROFILE_FILE):
        saved = pd.read_csv(PROFILE_FILE)
        st.session_state.profile = {
            "name": saved.iloc[0].get("昵称", ""),
            "time_slots": saved.iloc[0].get("学习时段", "").split("、"),
            "duration": saved.iloc[0].get("单次时长", "60分钟"),
            "style": saved.iloc[0].get("学习风格", ""),
            "weak_subjects": saved.iloc[0].get("薄弱科目", ""),
            "notes": saved.iloc[0].get("补充说明", ""),
            "setup_done": True
        }
    else:
        st.session_state.profile = {"setup_done": False}


# ============================================================
# 五、保存函数
# ============================================================
def save_ddl():
    st.session_state.df.to_csv(DDL_FILE, index=False, encoding="utf-8-sig")

def save_library():
    st.session_state.library.to_csv(LIBRARY_FILE, index=False, encoding="utf-8-sig")

def save_categories():
    pd.DataFrame({"分类": st.session_state.categories}).to_csv(CATEGORIES_FILE, index=False)


# ===== 【新增】保存学习画像 =====
def save_profile(profile_dict):
    """将学习画像持久化到 CSV"""
    df = pd.DataFrame([{
        "昵称": profile_dict.get("name", ""),
        "学习时段": "、".join(profile_dict.get("time_slots", [])),
        "单次时长": profile_dict.get("duration", "60分钟"),
        "学习风格": profile_dict.get("style", ""),
        "薄弱科目": profile_dict.get("weak_subjects", ""),
        "补充说明": profile_dict.get("notes", ""),
        "更新时间": datetime.now().strftime("%Y-%m-%d %H:%M")
    }])
    df.to_csv(PROFILE_FILE, index=False, encoding="utf-8-sig")


# ===== 【新增】生成个性化复习 Prompt =====
def generate_personalized_prompt(profile, tasks_df):
    """
    根据学习画像 + 当前未完成任务，生成一段结构化 Prompt。
    用户可以一键复制后粘贴到任意 AI 中获取定制化复习方案。
    """
    p = profile
    slots_str = "、".join(p.get("time_slots", ["未设置"]))

    # 构建任务摘要
    if tasks_df.empty:
        tasks_summary = "当前没有未完成的任务。"
    else:
        tasks_df = tasks_df.sort_values("截止日期")
        lines = []
        today = datetime.now().date()
        for _, row in tasks_df.iterrows():
            try:
                ddl = datetime.strptime(row["截止日期"], "%Y-%m-%d").date()
                days_left = (ddl - today).days
                urgency = "🔴 紧急" if days_left <= 2 else "🟡 较急" if days_left <= 5 else "🟢 从容"
            except:
                days_left = "?"
                urgency = "⚪ 未知"
            lines.append(
                f"  - {row['课程/科目']} | 截止: {row['截止日期']}（剩余{days_left}天）| "
                f"紧急度: {urgency} | 描述: {row['描述']}"
            )
        tasks_summary = "\n".join(lines)

    prompt = f"""你是一位专业的学习规划师，请根据以下学生画像和任务清单，生成一份详细的【一周个性化复习方案】。

━━━ 学生画像 ━━━
• 昵称：{p.get('name', '同学')}
• 每日可用学习时段：{slots_str}
• 每次学习时长偏好：{p.get('duration', '60分钟')}
• 学习风格：{p.get('style', '未填写')}
• 薄弱科目：{p.get('weak_subjects', '未填写')}
• 补充说明：{p.get('notes', '无')}

━━━ 当前未完成任务 ━━━
{tasks_summary}

━━━ 输出要求 ━━━
1. 按天（周一到周日）安排复习内容
2. 每天分配到具体的【可用学习时段】，每个时段给出明确的学习内容
3. 紧急任务（剩余≤2天）优先安排
4. 考虑薄弱科目的额外倾斜
5. 每天留出适当的休息和缓冲时间
6. 在方案末尾给出本周学习建议和注意事项
7. 用 Markdown 格式输出，包含表格或列表

请开始生成方案："""

    return prompt


# ============================================================
# 六、辅助工具函数
# ============================================================
def parse_flexible_date(date_str):
    if not date_str:
        return None
    date_str = date_str.strip().replace("。", "").replace("，", ",")
    for fmt in ("%Y-%m-%d", "%Y/%m/%d", "%m/%d"):
        try:
            dt = datetime.strptime(date_str, fmt)
            return dt.strftime("%Y-%m-%d")
        except ValueError:
            continue
    if date_parser:
        try:
            dt = date_parser.parse(date_str, fuzzy=True)
            return dt.strftime("%Y-%m-%d")
        except (ValueError, OverflowError):
            pass
    return None


def extract_text_from_file(uploaded_file):
    text = ""
    file_type = uploaded_file.type
    if file_type == "application/pdf":
        if PdfReader:
            reader = PdfReader(uploaded_file)
            for page in reader.pages:
                text += page.extract_text()
        else:
            st.error("请安装 pypdf")
    elif "word" in file_type or "document" in file_type:
        if Document:
            doc = Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        else:
            st.error("请安装 python-docx")
    elif "presentation" in file_type or "powerpoint" in file_type:
        if Presentation:
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            st.error("请安装 python-pptx")
    elif file_type.startswith("image/"):
        try:
            import easyocr, tempfile
            reader = easyocr.Reader(['ch_sim', 'en'], gpu=False)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp:
                tmp.write(uploaded_file.read())
                tmp_path = tmp.name
            result = reader.readtext(tmp_path, detail=0, paragraph=True)
            text = " ".join(result)
        except ImportError:
            st.error("请安装 easyocr 和 opencv-python-headless")
        except Exception as e:
            st.error(f"OCR失败：{e}")
    else:
        try:
            text = uploaded_file.read().decode("utf-8")
        except UnicodeDecodeError:
            try:
                uploaded_file.seek(0)
                text = uploaded_file.read().decode("gbk", errors="ignore")
            except Exception:
                st.error("无法解码文件")
    return text


# ============================================================
# 七、标题
# ============================================================
st.title("🎓 学习助手 Pro")


# ============================================================
# ===== 【新增】八、首次进入 — 学习画像配置引导 =====
# ============================================================
if not st.session_state.profile.get("setup_done", False):
    st.markdown("---")
    st.markdown("## 👋 欢迎使用学习助手 Pro！")
    st.markdown("首次使用，请花 1 分钟配置你的**个人学习画像**，系统将据此为你生成定制化复习方案。")
    st.markdown("")

    with st.form("profile_setup_form"):
        st.markdown("### ⏰ 设置你的学习时间")

        # 昵称
        name = st.text_input("你的昵称", placeholder="例如：小明")

        # 多选时间段
        st.markdown("**选择你每天可用于学习的时段**（可多选）：")
        slot_options = {
            "🌅 清晨 (6:00-8:00)": "清晨 6:00-8:00",
            "☀️ 上午 (8:00-12:00)": "上午 8:00-12:00",
            "🌤️ 下午 (13:00-17:00)": "下午 13:00-17:00",
            "🌆 傍晚 (17:00-19:00)": "傍晚 17:00-19:00",
            "🌙 晚间 (19:00-22:00)": "晚间 19:00-22:00",
            "🦉 深夜 (22:00-24:00)": "深夜 22:00-24:00",
        }

        col_a, col_b = st.columns(2)
        selected_slots = []
        for i, (label, value) in enumerate(slot_options.items()):
            with col_a if i % 2 == 0 else col_b:
                if st.checkbox(label, key=f"slot_{i}"):
                    selected_slots.append(value)

        st.markdown("")
        # 单次学习时长
        duration = st.selectbox(
            "每次学习时长偏好",
            ["30分钟", "45分钟", "60分钟", "90分钟", "120分钟"],
            index=2
        )

        # 学习风格
        style = st.selectbox(
            "你的学习风格",
            ["视觉型（喜欢图表、思维导图）",
             "听觉型（喜欢听课、讨论）",
             "阅读型（喜欢读教材、笔记）",
             "实践型（喜欢做题、实验）",
             "不确定"],
            index=4
        )

        # 薄弱科目
        weak = st.text_input("薄弱科目（可选）", placeholder="例如：高等数学、英语听力")

        # 补充说明
        notes = st.text_area("补充说明（可选）", placeholder="例如：周末白天有兼职，只有晚上能学习",
                             height=68)

        st.markdown("")
        submitted = st.form_submit_button("✅ 保存我的学习画像", use_container_width=True)

        if submitted:
            if not selected_slots:
                st.error("请至少选择一个学习时段")
            else:
                st.session_state.profile = {
                    "name": name if name else "同学",
                    "time_slots": selected_slots,
                    "duration": duration,
                    "style": style,
                    "weak_subjects": weak,
                    "notes": notes,
                    "setup_done": True
                }
                save_profile(st.session_state.profile)
                st.success("学习画像已保存！正在进入系统...")
                st.rerun()


# ============================================================
# ===== 【新增】九、学习画像展示卡片（配置完成后显示） =====
# ============================================================
if st.session_state.profile.get("setup_done", False):
    p = st.session_state.profile
    slots_html = "".join([f'<span class="tag">{s}</span>' for s in p.get("time_slots", [])])
    weak_text = p.get("weak_subjects", "") or "无"
    style_text = p.get("style", "未设置")

    st.markdown(f"""
    <div class="profile-card">
        <h4>📋 {p.get('name', '同学')}的学习画像</h4>
        <div>{slots_html}</div>
        <div class="info">
            ⏱ 单次时长: {p.get('duration', '60分钟')} &nbsp;|&nbsp;
            📖 风格: {style_text} &nbsp;|&nbsp;
            ⚠️ 薄弱: {weak_text}
        </div>
    </div>
    """, unsafe_allow_html=True)


# ============================================================
# 十、侧边栏
# ============================================================
with st.sidebar:
    st.header("⚙️ 设置")
    api_key = st.text_input("DeepSeek API Key", type="password", help="platform.deepseek.com 获取")

    st.divider()

    # ===== 【新增】学习画像管理入口 =====
    st.subheader("👤 学习画像")
    if st.session_state.profile.get("setup_done", False):
        if st.button("✏️ 修改学习画像"):
            st.session_state.profile["setup_done"] = False
            st.rerun()
    else:
        st.info("尚未配置，请在主页面完成设置")

    st.divider()

    st.subheader("🎨 主题")
    if st.button("切换暗黑模式"):
        if "dark" not in st.session_state:
            st.session_state.dark = False
        st.session_state.dark = not st.session_state.dark
        if st.session_state.dark:
            st.markdown("""
            <style>
            .stApp { background-color: #1e1e1e; color: #fff; }
            .stApp * { color: #eee; }
            .stButton button { background-color: #333; color: white; }
            </style>
            """, unsafe_allow_html=True)
        else:
            st.markdown("""
            <style>
            .stApp { background-color: #fff; color: #000; }
            </style>
            """, unsafe_allow_html=True)

    st.divider()

    st.subheader("💾 数据备份")
    json_ddl = st.session_state.df.to_json(orient="records", force_ascii=False)
    st.download_button("📥 备份DDL", data=json_ddl, file_name="backup_ddl.json")
    json_lib = st.session_state.library.to_json(orient="records", force_ascii=False)
    st.download_button("📥 备份资料库", data=json_lib, file_name="backup_library.json")
    uploaded_backup = st.file_uploader("恢复DDL备份", type=["json"], key="restore_ddl")
    if uploaded_backup:
        try:
            new_data = pd.DataFrame(json.loads(uploaded_backup.read()))
            if not new_data.empty:
                st.session_state.df = new_data
                save_ddl()
                st.success("恢复成功")
                st.rerun()
        except Exception as e:
            st.error(f"恢复失败：{e}")


# ============================================================
# 如果还没配置画像，只显示配置引导，不显示主功能
# ============================================================
if not st.session_state.profile.get("setup_done", False):
    st.stop()


# ============================================================
# 十一、智能输入区
# ============================================================
st.markdown("### ✨ 智能输入")
with st.container():
    input_col, btn_col = st.columns([5, 1])
    with input_col:
        smart_input = st.text_area("输入DDL或上传文件", placeholder="例如：下周一交高数作业，或上传课件",
                                   height=80, key="smart_input")
    with btn_col:
        st.write("")
        st.write("")
        send_btn = st.button("🚀 发送", use_container_width=True)

    uploaded_file = st.file_uploader("或点击上传文件 (PDF/Word/PPT/图片/TXT)", type=None,
                                     key="smart_file", label_visibility="collapsed")

    if send_btn:
        if uploaded_file is not None:
            content = extract_text_from_file(uploaded_file)
            if content:
                st.session_state["uploaded_content"] = content
                st.session_state["uploaded_filename"] = uploaded_file.name
                st.success(f"文件 '{uploaded_file.name}' 已读取，请到资料库选择分类保存")
            else:
                st.warning("文件内容提取失败")
        elif smart_input.strip():
            if not api_key:
                st.error("请先在侧边栏设置 API Key")
            else:
                with st.spinner("AI解析中..."):
                    try:
                        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
                        prompt = f"""提取学习任务信息。规则：
1. 如果是DDL，提取课程名称、截止日期(转为YYYY-MM-DD)、描述。
2. 如果不是DDL，请返回 {{"type": "other", "content": "原文"}}。
只返回JSON。
文本：{smart_input}"""
                        payload = {"model": "deepseek-chat", "messages": [{"role": "user", "content": prompt}],
                                   "temperature": 0.1}
                        response = requests.post("https://api.deepseek.com/v1/chat/completions",
                                                 headers=headers, json=payload, timeout=30)
                        if response.status_code == 200:
                            content = response.json()["choices"][0]["message"]["content"]
                            content = content.replace("```json", "").replace("```", "").strip()
                            parsed = json.loads(content)
                            if parsed.get("course") and parsed.get("deadline"):
                                st.session_state["parsed_course"] = parsed.get("course", "")
                                st.session_state["parsed_deadline"] = parsed.get("deadline", "")
                                st.session_state["parsed_notes"] = parsed.get("notes", "")
                                st.success("✅ 识别为DDL，已自动填入下方表单")
                            else:
                                st.info("识别为普通文本，已复制到剪贴板（可手动添加到资料库）")
                                st.session_state["smart_text"] = parsed.get("content", smart_input)
                        else:
                            st.error("AI解析失败")
                    except Exception as e:
                        st.error(f"出错：{e}")
        else:
            st.warning("请输入内容或上传文件")

st.divider()

# ============================================================
# 十二、主标签页
# ============================================================
tab_ddl, tab_lib = st.tabs(["📝 DDL管理（全功能）", "📚 资料库"])

# ==================== TAB 1: DDL 管理 ====================
with tab_ddl:

    # ----- 添加 DDL -----
    st.subheader("➕ 添加DDL")
    with st.form("ddl_form", clear_on_submit=True):
        col1, col2 = st.columns(2)
        with col1:
            course = st.text_input("课程/科目", value=st.session_state.get("parsed_course", ""))
            deadline_raw = st.text_input("截止日期 (支持多种格式)",
                                         value=st.session_state.get("parsed_deadline", ""))
        with col2:
            notes = st.text_input("描述", value=st.session_state.get("parsed_notes", ""))
            tag = st.text_input("标签", placeholder="作业/考试")
            repeat = st.selectbox("重复", ["无", "每周", "每月"])
        submitted = st.form_submit_button("💾 保存DDL")
        if submitted:
            deadline = parse_flexible_date(deadline_raw)
            if not course or not deadline:
                st.error("课程和截止日期必填")
            else:
                new_rows = []
                try:
                    start_date = datetime.strptime(deadline, "%Y-%m-%d")
                    dates = []
                    if repeat == "无":
                        dates = [start_date]
                    elif repeat == "每周":
                        dates = [start_date + timedelta(weeks=i) for i in range(4)]
                    elif repeat == "每月":
                        for i in range(3):
                            m = start_date.month + i
                            y = start_date.year + (m - 1) // 12
                            m = (m - 1) % 12 + 1
                            try:
                                dates.append(datetime(y, m, start_date.day))
                            except ValueError:
                                dates.append(datetime(y, m, 1))
                except ValueError:
                    st.error("日期格式错误")
                    st.stop()
                for dt in dates:
                    new_rows.append({
                        "课程/科目": course, "截止日期": dt.strftime("%Y-%m-%d"),
                        "描述": notes, "标签": tag if tag else "未分类",
                        "重复": repeat, "状态": "未完成",
                        "添加时间": datetime.now().strftime("%Y-%m-%d %H:%M")
                    })
                st.session_state.df = pd.concat([st.session_state.df, pd.DataFrame(new_rows)],
                                                ignore_index=True)
                save_ddl()
                for key in ["parsed_course", "parsed_deadline", "parsed_notes"]:
                    if key in st.session_state: del st.session_state[key]
                st.success(f"添加了 {len(new_rows)} 条DDL")
                st.rerun()

    # ----- 管理与搜索 -----
    st.subheader("🔍 管理与搜索")
    search = st.text_input("搜索DDL", key="ddl_search")
    filter_status = st.selectbox("状态筛选", ["全部", "未完成", "已完成"], key="ddl_status")
    df_display = st.session_state.df.copy()
    if search:
        df_display = df_display[df_display["课程/科目"].str.contains(search, na=False) | df_display["描述"].str.contains(search, na=False)]
    if filter_status != "全部":
        df_display = df_display[df_display["状态"] == filter_status]

    if not df_display.empty:
        edited = st.data_editor(
            df_display[["课程/科目", "截止日期", "描述", "标签", "状态"]],
            column_config={"状态": st.column_config.SelectboxColumn("状态", options=["未完成", "已完成"])},
            hide_index=True, key="ddl_edit"
        )
        if not edited.equals(df_display[["课程/科目", "截止日期", "描述", "标签", "状态"]]):
            for _, row in edited.iterrows():
                mask = (st.session_state.df["课程/科目"] == row["课程/科目"]) & (
                            st.session_state.df["截止日期"] == row["截止日期"])
                if mask.any():
                    st.session_state.df.loc[mask, "状态"] = row["状态"]
            save_ddl()
            st.success("状态已更新")
            st.rerun()

        st.divider()
        st.subheader("🗑️ 批量操作")
        df_del = df_display.copy()
        df_del["选择"] = False
        edited_del = st.data_editor(
            df_del[["选择", "课程/科目", "截止日期", "状态"]],
            column_config={"选择": st.column_config.CheckboxColumn("勾选")},
            hide_index=True, key="ddl_del"
        )
        col_del1, col_del2 = st.columns(2)
        with col_del1:
            if st.button("删除选中"):
                selected = edited_del[edited_del["选择"] == True]
                if not selected.empty:
                    for _, row in selected.iterrows():
                        mask = (st.session_state.df["课程/科目"] == row["课程/科目"]) & (
                                    st.session_state.df["截止日期"] == row["截止日期"])
                        if mask.any():
                            st.session_state.df = st.session_state.df.drop(st.session_state.df[mask].index)
                    save_ddl()
                    st.success(f"删除了 {len(selected)} 项")
                    st.rerun()
        with col_del2:
            if st.button("清除所有已完成"):
                before = len(st.session_state.df)
                st.session_state.df = st.session_state.df[st.session_state.df["状态"] != "已完成"]
                save_ddl()
                st.success(f"清除 {before - len(st.session_state.df)} 项")
                st.rerun()
    else:
        st.info("没有匹配的DDL")

    # ----- 数据分析 -----
    st.subheader("📊 数据分析")
    df_chart = st.session_state.df[st.session_state.df["状态"] != "已完成"].copy()
    if df_chart.empty:
        st.info("没有未完成的任务")
    else:
        try:
            df_chart["截止日期"] = pd.to_datetime(df_chart["截止日期"], errors='coerce')
            df_chart = df_chart.dropna(subset=["截止日期"])
            today = datetime.now()
            future_30 = today + timedelta(days=30)
            df_chart = df_chart[(df_chart["截止日期"] >= today) & (df_chart["截止日期"] <= future_30)]
            if not df_chart.empty:
                day_order = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
                day_names_cn = ["周一", "周二", "周三", "周四", "周五", "周六", "周日"]
                df_chart["星期"] = pd.Categorical(df_chart["截止日期"].dt.day_name(), categories=day_order)
                count_df = df_chart.groupby("星期").size().reset_index(name="任务数量")
                count_df["星期"] = count_df["星期"].map(dict(zip(day_order, day_names_cn)))
                fig = px.bar(count_df, x="星期", y="任务数量", title="未来30天 DDL 分布",
                             color="任务数量", color_continuous_scale="Reds", text="任务数量")
                fig.update_layout(height=350)
                st.plotly_chart(fig, use_container_width=True)
                st.subheader("🏷️ 标签分布")
                tag_df = df_chart["标签"].value_counts().reset_index()
                tag_df.columns = ["标签", "数量"]
                fig_pie = px.pie(tag_df, values="数量", names="标签", hole=0.3)
                st.plotly_chart(fig_pie, use_container_width=True)
            else:
                st.info("未来30天无任务")
        except Exception as e:
            st.warning(f"图表生成失败：{e}")

    # ----- 月视图 -----
    st.subheader("📆 月视图")
    if "cal_year" not in st.session_state:
        st.session_state.cal_year = datetime.now().year
        st.session_state.cal_month = datetime.now().month
    c1, c2, c3 = st.columns([1, 2, 1])
    with c1:
        if st.button("◀"):
            if st.session_state.cal_month == 1:
                st.session_state.cal_month = 12
                st.session_state.cal_year -= 1
            else:
                st.session_state.cal_month -= 1
            st.rerun()
    with c2:
        st.write(f"### {st.session_state.cal_year}年{st.session_state.cal_month}月")
    with c3:
        if st.button("▶"):
            if st.session_state.cal_month == 12:
                st.session_state.cal_month = 1
                st.session_state.cal_year += 1
            else:
                st.session_state.cal_month += 1
            st.rerun()
    year = st.session_state.cal_year
    month = st.session_state.cal_month
    first = datetime(year, month, 1)
    if month == 12:
        last = datetime(year + 1, 1, 1) - timedelta(days=1)
    else:
        last = datetime(year, month + 1, 1) - timedelta(days=1)
    start_week = first.weekday()
    total_days = last.day
    df_cal = st.session_state.df[st.session_state.df["状态"] != "已完成"].copy()
    tasks = {}
    if not df_cal.empty:
        df_cal["截止日期"] = pd.to_datetime(df_cal["截止日期"], errors='coerce')
        df_cal = df_cal.dropna(subset=["截止日期"])
        for _, row in df_cal.iterrows():
            if row["截止日期"].year == year and row["截止日期"].month == month:
                key = row["截止日期"].strftime("%Y-%m-%d")
                tasks.setdefault(key, []).append(row["课程/科目"])
    cal_dates = [None] * start_week + [datetime(year, month, d) for d in range(1, total_days + 1)]
    while len(cal_dates) < 42:
        cal_dates.append(None)
    html = "<div class='cal-grid'>"
    for w in ["一", "二", "三", "四", "五", "六", "日"]:
        html += f"<div class='cal-weekday'>{w}</div>"
    today_str = datetime.now().strftime("%Y-%m-%d")
    for d in cal_dates:
        if d is None:
            html += "<div class='cal-cell' style='background:transparent;'></div>"
        else:
            key = d.strftime("%Y-%m-%d")
            t = tasks.get(key, [])
            color = "#e74c3c" if len(t) >= 3 else "#f39c12" if len(t) >= 1 else "#2ecc71"
            bg = "#fff3cd" if key == today_str else "#f9f9f9"
            html += f"<div class='cal-cell' style='border-top-color:{color};background:{bg};'>"
            html += f"<div class='date'>{d.day}</div>"
            for task in t[:2]:
                html += f"<div style='font-size:9px;'>{task[:4]}</div>"
            if len(t) > 2:
                html += f"<div style='font-size:9px;color:#888;'>+{len(t) - 2}</div>"
            html += "</div>"
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)

    # ----- 导出与分享 -----
    st.subheader("📤 导出与分享")
    exp_col1, exp_col2 = st.columns(2)
    with exp_col1:
        if st.button("⬇️ 导出CSV"):
            csv = st.session_state.df.to_csv(index=False)
            b64 = base64.b64encode(csv.encode()).decode()
            st.markdown(f'<a href="data:file/csv;base64,{b64}" download="deadlines.csv">下载CSV</a>',
                        unsafe_allow_html=True)
        if st.button("⬇️ 导出日历(.ics)"):
            if Calendar is None:
                st.error("请安装 icalendar")
            else:
                try:
                    cal = Calendar()
                    cal.add('prodid', '-//DDL Pro//cn//')
                    cal.add('version', '2.0')
                    for _, row in st.session_state.df.iterrows():
                        if row["状态"] == "已完成": continue
                        event = Event()
                        event.add('summary', row["课程/科目"])
                        event.add('description', row["描述"])
                        date_obj = datetime.strptime(row["截止日期"], "%Y-%m-%d").date()
                        event.add('dtstart', date_obj)
                        event.add('dtend', date_obj)
                        cal.add_component(event)
                    ics_data = cal.to_ical()
                    b64 = base64.b64encode(ics_data).decode()
                    st.markdown(
                        f'<a href="data:text/calendar;base64,{b64}" download="deadlines.ics">下载.ics</a>',
                        unsafe_allow_html=True)
                except Exception as e:
                    st.error(f"导出失败：{e}")
    with exp_col2:
        if st.button("🔄 生成分享卡片"):
            df_share = st.session_state.df[st.session_state.df["状态"] != "已完成"].copy()
            if df_share.empty:
                st.warning("没有未完成任务")
            else:
                df_share = df_share.sort_values("截止日期")
                lines = ["# 📋 我的DDL清单\n"]
                for _, row in df_share.iterrows():
                    lines.append(f"- **{row['课程/科目']}** 截止: {row['截止日期']} | {row['描述']}")
                lines.append(f"\n> 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
                md = "\n".join(lines)
                st.text_area("复制以下内容", md, height=200)
                st.download_button("下载.md", md, file_name="DDL清单.md")

    # ============================================================
    # ===== 【新增】AI 复习规划（原功能 + 私人定制 Prompt） =====
    # ============================================================
    st.subheader("🧠 AI 复习规划")

    # 获取未完成任务（两个功能共用）
    df_future = st.session_state.df[st.session_state.df["状态"] != "已完成"].copy()

    # --- 按钮行：左=直接生成 | 右=生成私人定制 Prompt ---
    plan_col1, plan_col2 = st.columns(2)

    # 按钮 1：原有功能 —— 直接调用 API 生成复习计划
    with plan_col1:
        if st.button("📅 生成复习计划"):
            if df_future.empty or not api_key:
                st.warning("请确保有未完成任务并配置API Key")
            else:
                with st.spinner("生成中..."):
                    try:
                        tasks_text = df_future[["课程/科目", "截止日期", "描述"]].head(10).to_string()
                        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
                        prompt = f"基于以下任务：\n{tasks_text}\n生成未来一周复习优先级清单（Markdown列表）"
                        payload = {"model": "deepseek-chat", "messages": [{"role": "user", "content": prompt}],
                                   "temperature": 0.7}
                        response = requests.post("https://api.deepseek.com/v1/chat/completions",
                                                 headers=headers, json=payload, timeout=30)
                        if response.status_code == 200:
                            plan = response.json()["choices"][0]["message"]["content"]
                            st.markdown(plan)
                        else:
                            st.error("生成失败")
                    except Exception as e:
                        st.error(f"出错：{e}")

    # ===== 【新增】按钮 2：一键生成私人定制 Prompt =====
    with plan_col2:
        if st.button("🎯 生成私人定制方案 Prompt"):
            if df_future.empty:
                st.warning("没有未完成任务，请先添加 DDL")
            else:
                # 调用函数生成个性化 Prompt
                personal_prompt = generate_personalized_prompt(st.session_state.profile, df_future)
                st.session_state["personal_prompt"] = personal_prompt

    # ===== 【新增】展示私人定制 Prompt（可一键复制） =====
    if "personal_prompt" in st.session_state:
        st.markdown("---")
        st.markdown("#### 🎯 你的私人定制复习方案 Prompt")
        st.markdown("以下 Prompt 已结合你的**学习画像**（时间段、时长、风格、薄弱科目）和**当前任务清单**，"
                    "复制后粘贴到任意 AI（ChatGPT / DeepSeek / Kimi 等）即可获得专属复习方案。")

        # 使用 st.code 提供代码块样式（右上角自带复制按钮）
        st.code(st.session_state["personal_prompt"], language=None)

        # 额外提供下载按钮（备用方案）
        st.download_button(
            label="📥 下载 Prompt (.txt)",
            data=st.session_state["personal_prompt"],
            file_name="私人复习方案Prompt.txt",
            mime="text/plain"
        )

        # 一键清空
        if st.button("🗑️ 清除 Prompt"):
            del st.session_state["personal_prompt"]
            st.rerun()


# ==================== TAB 2: 资料库 ====================
with tab_lib:
    st.subheader("📁 分类管理")
    col_cat1, col_cat2 = st.columns([3, 1])
    with col_cat1:
        st.write("当前分类：", ", ".join(st.session_state.categories))
    with col_cat2:
        new_cat = st.text_input("新建分类", key="new_cat")
        if st.button("➕ 添加"):
            if new_cat.strip() and new_cat.strip() not in st.session_state.categories:
                st.session_state.categories.append(new_cat.strip())
                save_categories()
                st.success(f"已添加 '{new_cat.strip()}'")
                st.rerun()
            else:
                st.warning("分类已存在或为空")
    with st.expander("🗑️ 删除分类（仅当为空）"):
        del_cat = st.selectbox("选择分类", st.session_state.categories, key="del_cat")
        if st.button("确认删除"):
            if not st.session_state.library[st.session_state.library["分类"] == del_cat].empty:
                st.error(f"分类 '{del_cat}' 下还有文件")
            else:
                st.session_state.categories.remove(del_cat)
                save_categories()
                st.success("已删除")
                st.rerun()

    st.divider()
    st.subheader("📤 上传文件")
    with st.form("upload_lib"):
        uploaded_file = st.file_uploader("选择文件 (PDF/Word/PPT/图片/TXT)", type=None, key="lib_upload")
        cat_options = st.session_state.categories + ["新建分类..."]
        selected_cat = st.selectbox("选择分类", cat_options, key="lib_cat")
        if selected_cat == "新建分类...":
            new_cat_name = st.text_input("新分类名称")
        else:
            new_cat_name = None
        notes = st.text_area("备注（可选）")
        submitted = st.form_submit_button("📥 保存到资料库")
        if submitted and uploaded_file is not None:
            if new_cat_name and new_cat_name.strip():
                final_cat = new_cat_name.strip()
                if final_cat not in st.session_state.categories:
                    st.session_state.categories.append(final_cat)
                    save_categories()
            else:
                final_cat = selected_cat
            content = extract_text_from_file(uploaded_file)
            if content is None or content == "":
                st.error("内容提取失败")
            else:
                summary = content[:200] + ("..." if len(content) > 200 else "")
                new_row = {"文件名": uploaded_file.name, "分类": final_cat, "摘要": summary,
                           "上传时间": datetime.now().strftime("%Y-%m-%d %H:%M"), "内容": content}
                st.session_state.library = pd.concat([st.session_state.library, pd.DataFrame([new_row])],
                                                     ignore_index=True)
                save_library()
                st.success(f"文件 '{uploaded_file.name}' 已保存到 '{final_cat}'")
                st.rerun()

    st.divider()
    st.subheader("📂 资料浏览")
    filter_cat = st.selectbox("按分类筛选", ["全部"] + st.session_state.categories, key="lib_filter")
    search_lib = st.text_input("搜索", key="lib_search")
    df_lib = st.session_state.library.copy()
    if filter_cat != "全部":
        df_lib = df_lib[df_lib["分类"] == filter_cat]
    if search_lib:
        df_lib = df_lib[df_lib["文件名"].str.contains(search_lib, na=False) | df_lib["摘要"].str.contains(search_lib, na=False) | df_lib["内容"].str.contains(search_lib, na=False)]

    if df_lib.empty:
        st.info("暂无资料")
    else:
        for cat in st.session_state.categories:
            cat_df = df_lib[df_lib["分类"] == cat]
            if cat_df.empty: continue
            st.markdown(f"### 📁 {cat}")
            for idx, row in cat_df.iterrows():
                with st.container():
                    col1, col2, col3 = st.columns([3, 1, 1])
                    with col1:
                        st.markdown(f"**{row['文件名']}**")
                        st.caption(row['摘要'])
                    with col2:
                        if st.button("📄 查看", key=f"view_{idx}"):
                            st.text_area("全文", row['内容'], height=150)
                    with col3:
                        if st.button("🗑️ 删除", key=f"del_{idx}"):
                            st.session_state.library = st.session_state.library.drop(index=idx).reset_index(
                                drop=True)
                            save_library()
                            st.rerun()
            st.divider()

    if st.button("导出资料库CSV"):
        csv = st.session_state.library.to_csv(index=False)
        b64 = base64.b64encode(csv.encode()).decode()
        st.markdown(f'<a href="data:file/csv;base64,{b64}" download="library.csv">下载CSV</a>',
                    unsafe_allow_html=True)

st.caption("💡 全部功能：DDL管理（AI解析/图表/月视图/分享/复习）+ 资料库（分类/上传/搜索）+ 个性化学习画像")
